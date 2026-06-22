from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import os
import shutil
from typing import List, Optional
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
from google import genai
from datetime import datetime
import asyncio
from concurrent.futures import ThreadPoolExecutor
from dotenv import load_dotenv
from supabase import create_client, Client
import uuid
from io import BytesIO

# Load environment variables
load_dotenv()

app = FastAPI(title="File Upload & Notes API with Supabase", version="1.0.0")

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:5173",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:5173",
        "https://smart-notes-application.netlify.app",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configuration
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found in .env file")

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
if not SUPABASE_URL or not SUPABASE_KEY:
    raise ValueError("Supabase credentials not found in .env file")

MAX_FILE_SIZE = int(os.getenv("MAX_FILE_SIZE", 10 * 1024 * 1024))
ALLOWED_EXTENSIONS = set(
    os.getenv("ALLOWED_EXTENSIONS", ".pdf,.doc,.docx,.ppt,.pptx").split(",")
)

# Initialize Supabase client
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# Initialize Gemini
client = genai.Client(api_key=GOOGLE_API_KEY)

# Thread pool for text extraction
executor = ThreadPoolExecutor(max_workers=4)


def extract_text_from_pdf(file_content: bytes) -> str:
    try:
        text = ""
        pdf_reader = PyPDF2.PdfReader(BytesIO(file_content))
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF: {str(e)}")


def extract_text_from_docx(file_content: bytes) -> str:
    try:
        doc = docx.Document(BytesIO(file_content))
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading DOCX: {str(e)}")


def extract_text_from_pptx(file_content: bytes) -> str:
    try:
        lines = []
        prs = Presentation(BytesIO(file_content))
        for slide in prs.slides:
            seen = set()
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t and t not in seen:
                        seen.add(t)
                        lines.append(t)
        return "\n".join(lines)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX: {str(e)}")


def extract_text(file_content: bytes, file_extension: str) -> str:
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
    }

    if file_extension in [".doc", ".ppt"]:
        try:
            import textract

            text = textract.process(BytesIO(file_content)).decode("utf-8")
            return text.strip()
        except ImportError:
            raise HTTPException(
                status_code=400,
                detail="Please install textract for .doc/.ppt support: pip install textract",
            )
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error extracting text from {file_extension} file: {str(e)}",
            )

    if file_extension not in extractors:
        raise HTTPException(
            status_code=400, detail=f"No text extractor available for {file_extension}"
        )

    return extractors[file_extension](file_content)


async def generate_notes_with_gemini(text: str) -> str:
    try:
        if len(text) > 30000:
            text = text[:30000] + "..."

        prompt = f"""
Extract the key information from the following document and present it as structured notes.

Format your response EXACTLY like this, using these four section headers with no deviation:

OVERVIEW:
One or two sentences describing what this document is about and its purpose.

KEY POINTS:
• First key point
• Second key point
• Third key point
(list all major points, arguments, or topics covered)

DETAILS:
• Specific detail, fact, statistic, date, name, or finding
• Another specific detail
(include concrete facts, numbers, and specifics from the document)

CONCLUSIONS:
• Main takeaway or recommendation
• Any action items, next steps, or final judgments

Document text:
{text}
"""

        response = await asyncio.get_event_loop().run_in_executor(
            executor,
            lambda: client.models.generate_content(
                model="gemini-1.5-pro",
                contents=prompt,
            ),
        )

        return response.text.strip()
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating notes with Gemini: {str(e)}"
        )


def generate_simple_notes(text: str) -> str:
    """Fallback: build a basic bullet-note structure without AI."""
    if not text:
        return "OVERVIEW:\nNo text content found in the file.\n\nKEY POINTS:\n• N/A\n\nDETAILS:\n• N/A\n\nCONCLUSIONS:\n• N/A"

    text = " ".join(text.split())
    sentences = [s.strip() for s in text.split(". ") if s.strip()]

    overview = sentences[0] if sentences else "Document content extracted."
    key_points = sentences[1:5] if len(sentences) > 1 else []
    details = sentences[5:10] if len(sentences) > 5 else []
    conclusions = [sentences[-1]] if len(sentences) > 10 else []

    def fmt(items):
        return "\n".join(f"• {s}." for s in items) if items else "• N/A"

    return (
        f"OVERVIEW:\n{overview}.\n\n"
        f"KEY POINTS:\n{fmt(key_points)}\n\n"
        f"DETAILS:\n{fmt(details)}\n\n"
        f"CONCLUSIONS:\n{fmt(conclusions)}"
    )


async def upload_to_supabase_storage(file_content: bytes, filename: str) -> tuple:
    try:
        unique_filename = f"{uuid.uuid4()}_{filename}"
        file_path = f"documents/{unique_filename}"
        supabase.storage.from_("documents").upload(file_path, file_content)
        public_url = supabase.storage.from_("documents").get_public_url(file_path)
        return public_url, file_path
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error uploading to Supabase: {str(e)}"
        )


async def save_file_to_db(
    filename: str,
    original_filename: str,
    file_path: str,
    file_size: int,
    file_type: str,
    summary: str,
    text_preview: str,
    word_count: int,
    summary_method: str,
    user_id: Optional[str] = None,
) -> dict:
    try:
        data = {
            "filename": filename,
            "original_filename": original_filename,
            "file_path": file_path,
            "file_size": file_size,
            "file_type": file_type,
            "summary": summary,
            "text_preview": text_preview,
            "word_count": word_count,
            "summary_method": summary_method,
            "metadata": {"uploaded_at": datetime.now().isoformat()},
        }
        if user_id:
            data["user_id"] = user_id

        result = supabase.table("files").insert(data).execute()
        return result.data[0] if result.data else None
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error saving to database: {str(e)}"
        )


@app.api_route("/", methods=["GET", "HEAD"])
def read_root():
    return {
        "message": "File Upload & Notes API with Supabase",
        "endpoints": {
            "/upload": "POST - Upload a single file",
            "/upload/multiple": "POST - Upload multiple files",
            "/summary/{file_id}": "GET - Get file notes by ID",
            "/summaries": "GET - Get all notes",
            "/files": "GET - List all files",
            "/files/{file_id}": "DELETE - Delete a file",
            "/health": "GET - Health check",
        },
    }


@app.post("/upload/")
async def upload_file(file: UploadFile = File(...), user_id: Optional[str] = None):
    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    file_content = await file.read()
    if len(file_content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"File size exceeds {MAX_FILE_SIZE / 1024 / 1024}MB limit",
        )

    try:
        text_content = extract_text(file_content, file_extension)

        notes = None
        word_count = 0
        summary_method = "simple"

        if text_content:
            word_count = len(text_content.split())
            try:
                notes = await generate_notes_with_gemini(text_content)
                summary_method = "gemini"
            except Exception:
                notes = generate_simple_notes(text_content)
                summary_method = "simple"

        public_url, file_path = await upload_to_supabase_storage(
            file_content, file.filename
        )

        db_record = await save_file_to_db(
            filename=file_path.split("/")[-1],
            original_filename=file.filename,
            file_path=public_url,
            file_size=len(file_content),
            file_type=file_extension,
            summary=notes,
            text_preview=text_content[:500] if text_content else None,
            word_count=word_count,
            summary_method=summary_method,
            user_id=user_id,
        )

        return JSONResponse(
            status_code=200,
            content={
                "message": "File uploaded successfully",
                "file_id": db_record["id"] if db_record else None,
                "filename": file.filename,
                "file_path": public_url,
                "file_size": len(file_content),
                "summary": notes,
                "text_preview": text_content[:500] if text_content else None,
                "word_count": word_count,
                "summary_method": summary_method,
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")


@app.post("/upload/multiple/")
async def upload_multiple_files(
    files: List[UploadFile] = File(...), user_id: Optional[str] = None
):
    uploaded = []
    errors = []

    for file in files:
        try:
            file_extension = Path(file.filename).suffix.lower()
            if file_extension not in ALLOWED_EXTENSIONS:
                errors.append(
                    {
                        "filename": file.filename,
                        "error": f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
                    }
                )
                continue

            file_content = await file.read()
            if len(file_content) > MAX_FILE_SIZE:
                errors.append(
                    {
                        "filename": file.filename,
                        "error": f"File size exceeds {MAX_FILE_SIZE / 1024 / 1024}MB limit",
                    }
                )
                continue

            text_content = extract_text(file_content, file_extension)
            notes = None
            word_count = 0
            summary_method = "simple"

            if text_content:
                word_count = len(text_content.split())
                try:
                    notes = await generate_notes_with_gemini(text_content)
                    summary_method = "gemini"
                except Exception:
                    notes = generate_simple_notes(text_content)
                    summary_method = "simple"

            public_url, file_path = await upload_to_supabase_storage(
                file_content, file.filename
            )

            db_record = await save_file_to_db(
                filename=file_path.split("/")[-1],
                original_filename=file.filename,
                file_path=public_url,
                file_size=len(file_content),
                file_type=file_extension,
                summary=notes,
                text_preview=text_content[:500] if text_content else None,
                word_count=word_count,
                summary_method=summary_method,
                user_id=user_id,
            )

            uploaded.append(
                {
                    "filename": file.filename,
                    "file_id": db_record["id"] if db_record else None,
                    "file_path": public_url,
                    "file_size": len(file_content),
                    "summary": notes,
                    "word_count": word_count,
                    "summary_method": summary_method,
                }
            )
        except Exception as e:
            errors.append({"filename": file.filename, "error": str(e)})

    return JSONResponse(
        status_code=200,
        content={"uploaded": uploaded, "errors": errors if errors else None},
    )


@app.get("/files/")
async def list_files():
    try:
        result = (
            supabase.table("files").select("*").order("created_at", desc=True).execute()
        )
        return {"files": result.data}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error listing files: {str(e)}")


@app.get("/summary/{file_id}")
async def get_file_summary(file_id: str):
    try:
        result = supabase.table("files").select("*").eq("id", file_id).execute()
        if not result.data:
            raise HTTPException(status_code=404, detail="File not found")
        return JSONResponse(status_code=200, content=result.data[0])
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/summaries/")
async def get_all_summaries():
    try:
        result = (
            supabase.table("files")
            .select(
                "id, original_filename, summary, word_count, summary_method, created_at"
            )
            .order("created_at", desc=True)
            .execute()
        )
        return JSONResponse(
            status_code=200,
            content={"total_files": len(result.data), "summaries": result.data},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.delete("/files/{file_id}")
async def delete_file(file_id: str):
    try:
        record = supabase.table("files").select("*").eq("id", file_id).execute()
        if not record.data:
            raise HTTPException(status_code=404, detail="File not found")

        file_path = record.data[0]["file_path"]
        storage_path = "documents/" + file_path.split("/documents/")[-1]
        supabase.storage.from_("documents").remove([storage_path])
        supabase.table("files").delete().eq("id", file_id).execute()

        return {"message": f"File '{file_id}' deleted successfully"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting file: {str(e)}")


@app.get("/health")
async def health_check():
    gemini_status = "not_connected"
    supabase_status = "not_connected"

    try:
        test = client.models.generate_content(
            model="gemini-1.5-pro", contents="Respond with OK"
        )
        gemini_status = "connected" if test.text else "error"
    except Exception:
        pass

    try:
        supabase.table("files").select("count", count="exact").execute()
        supabase_status = "connected"
    except Exception:
        pass

    return {
        "status": "healthy",
        "gemini": gemini_status,
        "supabase": supabase_status,
    }


if __name__ == "__main__":
    import uvicorn

    host = os.getenv("HOST", "0.0.0.0")
    port = int(os.getenv("PORT", 8000))
    reload = os.getenv("RELOAD", "True").lower() == "true"
    uvicorn.run("main:app", host=host, port=port, reload=reload)
