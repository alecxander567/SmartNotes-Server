from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
import os
import shutil
from typing import List, Optional
from pathlib import Path
import PyPDF2
import docx
from pptx import Presentation
import google.generativeai as genai
from datetime import datetime
import asyncio
from concurrent.futures import ThreadPoolExecutor

app = FastAPI()

# Configuration
UPLOAD_DIR = "uploads"
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
ALLOWED_EXTENSIONS = {".pdf", ".doc", ".docx", ".ppt", ".pptx"}

# Google Gemini Configuration
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY environment variable is not set")

genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel(
    "gemini-1.5-pro"
)  # or 'gemini-1.5-flash' for faster responses

# Create upload directory
Path(UPLOAD_DIR).mkdir(parents=True, exist_ok=True)

# Store file data in memory (use database in production)
file_data = {}

# Thread pool for text extraction
executor = ThreadPoolExecutor(max_workers=4)


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF file."""
    try:
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF: {str(e)}")


def extract_text_from_docx(file_path: str) -> str:
    """Extract text from DOCX file."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading DOCX: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file."""
    try:
        text = ""
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        text += shape.text + "\n"
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text:
                            text += paragraph.text + "\n"
        return text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PPTX: {str(e)}")


def extract_text(file_path: str, file_extension: str) -> str:
    """Extract text based on file type."""
    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
    }

    if file_extension in [".doc", ".ppt"]:
        # Try to use textract for older formats
        try:
            import textract

            text = textract.process(file_path).decode("utf-8")
            return text.strip()
        except ImportError:
            raise HTTPException(
                status_code=400,
                detail=f"Please install textract for .doc/.ppt support: pip install textract",
            )
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"Error extracting text from {file_extension} file: {str(e)}",
            )

    if file_extension not in extractors:
        raise HTTPException(
            status_code=400, detail=f"No text extractor available for {file_extension}"
        )

    return extractors[file_extension](file_path)


async def generate_summary_with_gemini(text: str, max_length: int = 200) -> str:
    """Generate summary using Google Gemini API."""
    try:
        # Truncate text if too long (Gemini has a context window limit)
        if len(text) > 30000:
            text = text[:30000] + "..."

        prompt = f"""
        Please provide a concise and informative summary of the following document in {max_length} words or less.
        
        The summary should:
        1. Capture the main topic and purpose
        2. Highlight key points and findings
        3. Include important conclusions or recommendations
        4. Be written in clear, professional language
        
        Document text:
        {text}
        
        Summary:
        """

        response = await asyncio.get_event_loop().run_in_executor(
            executor, model.generate_content, prompt
        )

        summary = response.text.strip()
        return summary
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating summary with Gemini: {str(e)}"
        )


def generate_simple_summary(text: str, max_length: int = 200) -> str:
    """Generate a simple extractive summary (fallback without AI)."""
    if not text:
        return "No text content found in the file."

    # Remove extra whitespace and split into sentences
    text = " ".join(text.split())
    sentences = text.split(". ")

    # If text is short, return first few sentences
    if len(text) <= max_length:
        return text

    if len(sentences) <= 3:
        return text[:max_length] + "..."

    # Get first and last few sentences as a simple summary
    summary_sentences = sentences[:2] + sentences[-2:]
    summary = ". ".join(summary_sentences)

    # Truncate if needed
    if len(summary) > max_length:
        summary = summary[:max_length] + "..."

    return summary


def get_file_metadata(file_path: str, filename: str) -> dict:
    """Get file metadata."""
    stat = os.stat(file_path)
    return {
        "filename": filename,
        "size_bytes": stat.st_size,
        "size_mb": round(stat.st_size / (1024 * 1024), 2),
        "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "extension": Path(filename).suffix.lower(),
    }


@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    """Upload a single file and generate summary."""
    # Validate file extension
    file_extension = Path(file.filename).suffix.lower()
    if file_extension not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    # Check file size
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)

    if file_size > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"File size exceeds {MAX_FILE_SIZE / 1024 / 1024}MB limit",
        )

    # Generate unique filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    unique_filename = f"{timestamp}_{file.filename}"
    file_path = os.path.join(UPLOAD_DIR, unique_filename)

    try:
        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Extract text
        text_content = extract_text(file_path, file_extension)

        # Generate summary using Gemini
        summary = None
        word_count = 0
        if text_content:
            word_count = len(text_content.split())
            try:
                summary = await generate_summary_with_gemini(text_content)
            except Exception as e:
                # Fallback to simple summary if Gemini fails
                summary = generate_simple_summary(text_content)

        # Store file data
        file_data[unique_filename] = {
            "original_filename": file.filename,
            "summary": summary,
            "text_preview": text_content[:500] if text_content else None,
            "word_count": word_count,
            "metadata": get_file_metadata(file_path, unique_filename),
            "uploaded_at": datetime.now().isoformat(),
            "summary_method": "gemini" if summary else "simple",
        }

        return JSONResponse(
            status_code=200,
            content={
                "message": "File uploaded successfully",
                "filename": file.filename,
                "saved_as": unique_filename,
                "file_size": file_size,
                "summary": summary,
                "word_count": word_count,
                "summary_method": "gemini" if summary else "simple",
            },
        )
    except HTTPException:
        # Clean up if error occurs
        if os.path.exists(file_path):
            os.remove(file_path)
        raise
    except Exception as e:
        # Clean up if error occurs
        if os.path.exists(file_path):
            os.remove(file_path)
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")


@app.post("/upload/multiple/")
async def upload_multiple_files(files: List[UploadFile] = File(...)):
    """Upload multiple files."""
    uploaded = []
    errors = []

    for file in files:
        try:
            # Validate file extension
            file_extension = Path(file.filename).suffix.lower()
            if file_extension not in ALLOWED_EXTENSIONS:
                errors.append(
                    {
                        "filename": file.filename,
                        "error": f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
                    }
                )
                continue

            # Check file size
            file.file.seek(0, 2)
            file_size = file.file.tell()
            file.file.seek(0)

            if file_size > MAX_FILE_SIZE:
                errors.append(
                    {
                        "filename": file.filename,
                        "error": f"File size exceeds {MAX_FILE_SIZE / 1024 / 1024}MB limit",
                    }
                )
                continue

            # Save file
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            unique_filename = f"{timestamp}_{file.filename}"
            file_path = os.path.join(UPLOAD_DIR, unique_filename)

            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)

            # Extract text and generate summary
            text_content = extract_text(file_path, file_extension)
            summary = None
            word_count = 0

            if text_content:
                word_count = len(text_content.split())
                try:
                    summary = await generate_summary_with_gemini(text_content)
                except:
                    summary = generate_simple_summary(text_content)

            # Store file data
            file_data[unique_filename] = {
                "original_filename": file.filename,
                "summary": summary,
                "text_preview": text_content[:500] if text_content else None,
                "word_count": word_count,
                "metadata": get_file_metadata(file_path, unique_filename),
                "uploaded_at": datetime.now().isoformat(),
                "summary_method": "gemini" if summary else "simple",
            }

            uploaded.append(
                {
                    "filename": file.filename,
                    "saved_as": unique_filename,
                    "file_size": file_size,
                    "summary": summary,
                    "word_count": word_count,
                }
            )
        except Exception as e:
            errors.append({"filename": file.filename, "error": str(e)})

    return JSONResponse(
        status_code=200,
        content={"uploaded": uploaded, "errors": errors if errors else None},
    )


@app.get("/summary/{filename}")
async def get_file_summary(filename: str):
    """Get summary of a specific uploaded file."""
    if filename not in file_data:
        # Check if file exists but not in memory
        file_path = os.path.join(UPLOAD_DIR, filename)
        if os.path.exists(file_path):
            # Try to regenerate summary
            file_extension = Path(filename).suffix.lower()
            try:
                text_content = extract_text(file_path, file_extension)
                if text_content:
                    summary = await generate_summary_with_gemini(text_content)
                    file_data[filename] = {
                        "original_filename": filename,
                        "summary": summary,
                        "text_preview": text_content[:500] if text_content else None,
                        "word_count": len(text_content.split()) if text_content else 0,
                        "metadata": get_file_metadata(file_path, filename),
                        "uploaded_at": datetime.now().isoformat(),
                        "summary_method": "gemini",
                    }
                    return JSONResponse(status_code=200, content=file_data[filename])
                else:
                    file_data[filename] = {
                        "original_filename": filename,
                        "summary": "No text content found in file",
                        "text_preview": None,
                        "word_count": 0,
                        "metadata": get_file_metadata(file_path, filename),
                        "uploaded_at": datetime.now().isoformat(),
                        "summary_method": "none",
                    }
            except:
                pass

        raise HTTPException(status_code=404, detail="File not found")

    return JSONResponse(status_code=200, content=file_data[filename])


@app.get("/summaries/")
async def get_all_summaries():
    """Get summaries of all uploaded files."""
    summaries = {}
    for filename, data in file_data.items():
        summaries[filename] = {
            "original_filename": data["original_filename"],
            "summary": data["summary"],
            "word_count": data["word_count"],
            "summary_method": data.get("summary_method", "unknown"),
            "uploaded_at": data.get("uploaded_at"),
        }

    return JSONResponse(
        status_code=200, content={"total_files": len(summaries), "summaries": summaries}
    )


@app.get("/files/")
async def list_files():
    """List all uploaded files."""
    try:
        files = []
        for filename in os.listdir(UPLOAD_DIR):
            file_path = os.path.join(UPLOAD_DIR, filename)
            if os.path.isfile(file_path):
                file_info = {
                    "filename": filename,
                    "size": os.path.getsize(file_path),
                    "modified": datetime.fromtimestamp(
                        os.path.getmtime(file_path)
                    ).isoformat(),
                }
                # Add summary if available
                if filename in file_data:
                    file_info["summary"] = file_data[filename].get("summary")
                    file_info["word_count"] = file_data[filename].get("word_count")
                    file_info["summary_method"] = file_data[filename].get(
                        "summary_method"
                    )
                files.append(file_info)
        return {"files": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error listing files: {str(e)}")


@app.delete("/files/{filename}")
async def delete_file(filename: str):
    """Delete a specific file and its data."""
    try:
        file_path = os.path.join(UPLOAD_DIR, filename)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")

        os.remove(file_path)

        # Remove from memory
        if filename in file_data:
            del file_data[filename]

        return {"message": f"File '{filename}' deleted successfully"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting file: {str(e)}")


@app.post("/generate-summary/{filename}")
async def generate_summary_for_file(filename: str):
    """Generate or regenerate summary for a specific file using Gemini."""
    file_path = os.path.join(UPLOAD_DIR, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    file_extension = Path(filename).suffix.lower()
    try:
        text_content = extract_text(file_path, file_extension)
        if not text_content:
            return JSONResponse(
                status_code=200,
                content={
                    "message": "No text content found in the file",
                    "filename": filename,
                },
            )

        # Generate summary using Gemini
        summary = await generate_summary_with_gemini(text_content)

        # Update stored data
        file_data[filename] = {
            "original_filename": filename,
            "summary": summary,
            "text_preview": text_content[:500] if text_content else None,
            "word_count": len(text_content.split()) if text_content else 0,
            "metadata": get_file_metadata(file_path, filename),
            "uploaded_at": datetime.now().isoformat(),
            "summary_method": "gemini",
        }

        return JSONResponse(
            status_code=200,
            content={
                "message": "Summary generated successfully using Gemini",
                "filename": filename,
                "summary": summary,
                "word_count": len(text_content.split()),
            },
        )
    except Exception as e:
        raise HTTPException(
            status_code=500, detail=f"Error generating summary: {str(e)}"
        )


# Health check endpoint
@app.get("/health")
async def health_check():
    """Check if the service is running and Gemini is configured."""
    try:
        # Test Gemini connection
        test_response = model.generate_content("Respond with 'OK'")
        gemini_status = "connected" if test_response.text else "error"
    except:
        gemini_status = "not_connected"

    return {
        "status": "healthy",
        "gemini": gemini_status,
        "uploaded_files": len(file_data),
    }
